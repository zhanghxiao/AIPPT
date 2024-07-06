import os
import requests
import json
from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from dotenv import load_dotenv
from io import BytesIO
from PIL import Image

app = Flask(__name__)

load_dotenv()

OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
API_REQUEST_URL = os.getenv('API_REQUEST_URL')
MINDMAP_API_URL = os.getenv('MINDMAP_API_URL')
MINDMAP_API_KEY = os.getenv('MINDMAP_API_KEY')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    data = request.json
    topic = data['topic']
    reference = data.get('reference', '')

    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {OPENAI_API_KEY}'
    }
    prompt = (
        f'请为主题"{topic}"生成一个详细的PPT内容，包括封面、目录和5个以上的内容页面,以及需要2张以上图片，封面不添加图片，必须有目录。内容页面必须有相关的内容。'
        '在适当的位置，您可以添加一个思维导图来总结或展示关键信息,整个PPT最多只能有两个思维导图。'
        '使用以下格式和标识符：\n'
        '1. 使用"[SLIDE]"作为每页幻灯片的开始标记。\n'
        '2. 使用"[TITLE]"标记主标题。\n'
        '3. 使用"[SUBTITLE]"标记副标题。\n'
        '4. 使用"[CONTENT]"标记普通内容，每个要点使用"-"开始。\n'
        '5. 使用"[IMAGE]"标记图片位置，后面跟图片描述。\n'
        '6. 使用"[MINDMAP]"标记思维导图位置，格式为"[MINDMAP]思维导图标题|思维导图参考内容"。\n'
        '例如：\n'
        '[SLIDE]\n'
        '[TITLE]PPT主题\n'
        '[SUBTITLE]副标题\n'
        '[CONTENT]\n'
        '- 内容要点1\n'
        '- 内容要点2\n'
        '[IMAGE]这里是图片描述\n'
        '[MINDMAP]主题概览|包含主要概念和关键点\n'
    )
    if reference:
        prompt += f' 参考内容或者修改意见: {reference}'

    data = {
        'model': 'gpt-3.5-turbo',
        'messages': [
            {'role': 'system', 'content': prompt},
            {'role': 'user', 'content': '请生成完整的PPT内容。'}
        ]
    }

    response = requests.post(API_REQUEST_URL, headers=headers, json=data)
    response_data = response.json()
    ppt_content = response_data['choices'][0]['message']['content']

    # 解析内容，提取思维导图信息
    mindmap_info = []
    slides = ppt_content.split('[SLIDE]')
    for slide in slides:
        if '[MINDMAP]' in slide:
            mindmap_line = [line for line in slide.split('\n') if line.startswith('[MINDMAP]')][0]
            mindmap_content = mindmap_line.replace('[MINDMAP]', '').strip().split('|')
            if len(mindmap_content) >= 2:
                mindmap_title = mindmap_content[0]
                mindmap_reference = mindmap_content[1]
                mindmap_url = generate_mindmap(mindmap_title, mindmap_reference)
                if mindmap_url:
                    mindmap_info.append({
                        'title': mindmap_title,
                        'url': mindmap_url
                    })

    ppt = create_ppt(ppt_content, mindmap_info)
    os.makedirs('static', exist_ok=True)
    ppt_path = os.path.join('static', 'generated_ppt.pptx')
    ppt.save(ppt_path)

    return jsonify({
        "content": ppt_content,
        "file_path": ppt_path,
        "mindmaps": mindmap_info
    })


@app.route('/update_ppt', methods=['POST'])
def update_ppt():
    content = request.json['content']

    # 解析内容，提取思维导图信息
    mindmap_info = []
    slides = content.split('[SLIDE]')
    for slide in slides:
        if '[MINDMAP]' in slide:
            mindmap_line = [line for line in slide.split('\n') if line.startswith('[MINDMAP]')][0]
            mindmap_content = mindmap_line.replace('[MINDMAP]', '').strip().split('|')
            if len(mindmap_content) >= 2:
                mindmap_title = mindmap_content[0]
                mindmap_reference = mindmap_content[1]
                mindmap_url = generate_mindmap(mindmap_title, mindmap_reference)
                if mindmap_url:
                    mindmap_info.append({
                        'title': mindmap_title,
                        'url': mindmap_url
                    })

    ppt = create_ppt(content, mindmap_info)
    ppt_path = os.path.join('static', 'generated_ppt.pptx')
    ppt.save(ppt_path)
    return jsonify({"success": True, "mindmaps": mindmap_info})

@app.route('/download_ppt')
def download_ppt():
    ppt_path = os.path.join('static', 'generated_ppt.pptx')
    if os.path.exists(ppt_path):
        return send_file(ppt_path, as_attachment=True)
    else:
        return "PPT file not found", 404

def add_image_placeholder(slide, description):
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(2.5)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)

    tf = shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = f"[图片: {description}]"
    p.alignment = PP_ALIGN.CENTER

    return shape


def create_ppt(content, mindmap_info):
    ppt = Presentation()
    slides = content.split('[SLIDE]')

    for i, slide_content in enumerate(slides[1:]):
        if i == 0:  # Cover slide
            slide = ppt.slides.add_slide(ppt.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]

        lines = slide_content.strip().split('\n')

        for line in lines:
            if line.startswith('[TITLE]'):
                title.text = line.replace('[TITLE]', '').strip()
            elif line.startswith('[SUBTITLE]'):
                if i == 0:
                    subtitle.text = line.replace('[SUBTITLE]', '').strip()
                else:
                    content.text += line.replace('[SUBTITLE]', '').strip() + '\n'
            elif line.startswith('[CONTENT]'):
                if i != 0:
                    content.text = ''
            elif line.startswith('-'):
                if i != 0:
                    p = content.text_frame.add_paragraph()
                    p.text = line.strip()
                    p.level = 0
            elif line.startswith('[IMAGE]'):
                description = line.replace('[IMAGE]', '').strip()
                add_image_placeholder(slide, description)
            elif line.startswith('[MINDMAP]'):
                mindmap_content = line.replace('[MINDMAP]', '').strip().split('|')
                if len(mindmap_content) >= 2:
                    mindmap_title = mindmap_content[0]
                    mindmap_reference = mindmap_content[1]
                    # 为思维导图创建新的幻灯片
                    mindmap_slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # 使用空白布局

                    # 检查是否存在标题占位符，如果不存在，手动添加文本框
                    if mindmap_slide.shapes.title:
                        mindmap_title_shape = mindmap_slide.shapes.title
                        mindmap_title_shape.text = "思维导图: " + mindmap_title
                    else:
                        left = Inches(0.5)
                        top = Inches(0.5)
                        width = Inches(9)
                        height = Inches(1)
                        textbox = mindmap_slide.shapes.add_textbox(left, top, width, height)
                        textbox.text_frame.text = "思维导图: " + mindmap_title

                    # 查找对应的思维导图 URL
                    mindmap_url = next((item['url'] for item in mindmap_info if item['title'] == mindmap_title), None)

                    if mindmap_url:
                        add_mindmap_to_slide(mindmap_slide, mindmap_url)
                    else:
                        add_image_placeholder(mindmap_slide, "思维导图生成失败")

        # Set font properties
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if i == 0 and shape == title:  # Cover slide title
                        run.font.name = '微软雅黑'
                        run.font.size = Pt(44)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    elif i == 0 and shape == subtitle:  # Cover slide subtitle
                        run.font.name = '微软雅黑'
                        run.font.size = Pt(32)
                        run.font.color.rgb = RGBColor(89, 89, 89)
                    else:  # Other slides
                        run.font.name = '微软雅黑'
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(0, 0, 0)

    return ppt

def generate_mindmap(title, reference):
    # 首先，调用 OpenAI API 生成思维导图内容
    mindmap_prompt = f"我想使用代码创建一个《{title}》的思维导图，其中包含多个主题和子主题，以及叶子节点。请你提供一些Markdown格式的文本。在Markdown格式中，首行必须是@startmindmap，结尾必须是@endmindmap，* 表示中央主题， ** 表示主要主题，*** 表示子主题，**** 表示叶子节点。请参照以上格式，在markdown代码块中帮我创建一个有效的思维导图。下面是参考内容：{reference}"

    openai_response = requests.post(
        API_REQUEST_URL,
        headers={
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {OPENAI_API_KEY}'
        },
        json={
            'model': 'gpt-3.5-turbo',
            'messages': [
                {'role': 'system', 'content': mindmap_prompt},
                {'role': 'user', 'content': '请生成思维导图内容。'}
            ]
        }
    )

    if openai_response.status_code != 200:
        return None

    mindmap_content = openai_response.json()['choices'][0]['message']['content']

    # 提取 @startmindmap 和 @endmindmap 之间的内容
    start_index = mindmap_content.find('@startmindmap')
    end_index = mindmap_content.find('@endmindmap')
    if start_index != -1 and end_index != -1:
        mindmap_content = mindmap_content[start_index:end_index + 11]
    else:
        return None

    # 调用图像生成服务 API 生成思维导图图像
    mindmap_response = requests.post(
        MINDMAP_API_URL,
        headers={
            'Authorization': f'Bearer {MINDMAP_API_KEY}',
            'Content-Type': 'application/json'
        },
        json={
            'model': 'dall-e-3',
            'prompt': mindmap_content,
            'n': 1,
            'size': '1024x1024'
        }
    )

    if mindmap_response.status_code != 200:
        return None

    mindmap_url = mindmap_response.json()['data'][0]['url']
    return mindmap_url

def add_mindmap_to_slide(slide, image_url):
    response = requests.get(image_url)
    image_data = BytesIO(response.content)

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5)

    slide.shapes.add_picture(image_data, left, top, width, height)

if __name__ == '__main__':
    app.run(debug=True)
